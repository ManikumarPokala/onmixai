"""PPTX parser (python-pptx): slide text + speaker notes, one block per slide."""

import io

from pptx import Presentation

from src.knowledge.ingest_errors import ParserError
from src.knowledge.parsing.base import ContentBlock, ParsedDocument


class PptxParser:
    def parse(self, data: bytes, *, max_pages: int) -> ParsedDocument:
        """Time: O(slides + shapes). Slides count as pages for the page cap."""
        try:
            presentation = Presentation(io.BytesIO(data))
        except Exception as exc:
            raise ParserError("file is not a readable PPTX") from exc

        slides = list(presentation.slides)
        if len(slides) > max_pages:
            raise ParserError(f"PPTX exceeds the page limit of {max_pages}")

        blocks: list[ContentBlock] = []
        for number, slide in enumerate(slides, start=1):
            texts = [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]
            body = "\n".join(text for text in texts if text.strip())
            if body:
                blocks.append(ContentBlock(text=body, ref={"slide": number}))
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text
                if notes.strip():
                    blocks.append(ContentBlock(text=notes, ref={"slide": number, "kind": "notes"}))
        return ParsedDocument(blocks=tuple(blocks), page_count=len(slides))
