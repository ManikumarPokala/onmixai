"""Deterministic generators for valid + broken document fixtures.

Checked-in code (not opaque committed binaries) so every fixture is reproducible.
Used by the parser tests and the Task 10 broken-corpus drill.
"""

import io

import pymupdf
from docx import Document as DocxDocument
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches

# --- valid fixtures ---


def valid_txt_utf8() -> bytes:
    return b"First paragraph.\n\nSecond paragraph with more words."


def valid_txt_legacy() -> bytes:
    # Accented text in latin-1 is invalid UTF-8 but resolvable by charset-normalizer.
    return ("Café crème, déjà vu. " * 12 + "Naïve façade, fiancé.").encode("latin-1")


def valid_pdf(pages: int = 2) -> bytes:
    document = pymupdf.open()
    try:
        for index in range(pages):
            page = document.new_page()
            page.insert_text((72, 72), f"Page {index + 1} body text with several words.")
        return bytes(document.tobytes())
    finally:
        document.close()


def scanned_pdf() -> bytes:
    """A one-page PDF whose page is an image with no text layer (forces OCR)."""
    document = pymupdf.open()
    try:
        page = document.new_page()
        pixmap = pymupdf.Pixmap(pymupdf.csRGB, pymupdf.IRect(0, 0, 200, 200))
        pixmap.clear_with(255)
        page.insert_image(page.rect, pixmap=pixmap)
        return bytes(document.tobytes())
    finally:
        document.close()


def valid_docx() -> bytes:
    document = DocxDocument()
    document.add_paragraph("Document heading paragraph.")
    document.add_paragraph("Another paragraph of prose content.")
    table = document.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "H1"
    table.rows[0].cells[1].text = "H2"
    table.rows[1].cells[0].text = "v1"
    table.rows[1].cells[1].text = "v2"
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def valid_pptx() -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "Slide one title and body."
    slide.notes_slide.notes_text_frame.text = "Speaker notes for slide one."
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def valid_xlsx() -> bytes:
    workbook = Workbook()
    sheet = workbook.active
    assert sheet is not None
    sheet.title = "Q1"
    sheet.append(["Region", "Revenue"])
    sheet.append(["EMEA", 100])
    sheet.append(["APAC", 200])
    buffer = io.BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


# --- broken fixtures ---


def zero_byte() -> bytes:
    return b""


def truncated_pdf() -> bytes:
    full = valid_pdf(1)
    return full[: len(full) // 2]


def password_pdf() -> bytes:
    document = pymupdf.open()
    try:
        document.new_page().insert_text((72, 72), "secret")
        return bytes(
            document.tobytes(
                encryption=pymupdf.PDF_ENCRYPT_AES_256, owner_pw="owner", user_pw="user"
            )
        )
    finally:
        document.close()


def png_as_pdf() -> bytes:
    """A real PNG mislabeled as a PDF."""
    pixmap = pymupdf.Pixmap(pymupdf.csRGB, pymupdf.IRect(0, 0, 50, 50))
    pixmap.clear_with(128)
    return bytes(pixmap.tobytes("png"))


def corrupt_docx() -> bytes:
    full = valid_docx()
    return full[: len(full) // 2]


def corrupt_xlsx() -> bytes:
    full = valid_xlsx()
    return full[: len(full) // 2]


def garbage_txt() -> bytes:
    """Deterministic non-text bytes — no encoding decodes it confidently."""
    return bytes((index * 37) % 256 for index in range(2048))
